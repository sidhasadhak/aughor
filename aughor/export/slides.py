"""report_json (parsed to an ExportDoc) → a PowerPoint deck, via python-pptx.

Layout strategy: a title slide, then one "section" slide per heading that
collects the prose/bullets/findings under it, with charts and tables promoted to
their own full-bleed slides (they read better large). Walks the SAME ExportDoc
blocks the PDF does — no report-shape logic here.
"""
from __future__ import annotations

import io
import re
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .document import Block, ExportDoc

_INDIGO = RGBColor(0x4F, 0x46, 0xE5)
_INK = RGBColor(0x18, 0x18, 0x1B)
_BODY = RGBColor(0x3F, 0x3F, 0x46)
_MUTED = RGBColor(0x71, 0x71, 0x7A)
_BG = RGBColor(0xF4, 0xF4, 0xF5)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

_W = Inches(13.333)
_H = Inches(7.5)
_MARGIN = Inches(0.7)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _title_bar(slide, text: str):
    box = slide.shapes.add_textbox(_MARGIN, Inches(0.45), _W - _MARGIN * 2, Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = _INDIGO
    # accent underline
    line = slide.shapes.add_shape(1, _MARGIN, Inches(1.35), Inches(1.1), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = _INDIGO
    line.line.fill.background()
    return slide


_BOLD_RE = re.compile(r"\*\*(.+?)\*\*")


def _para(tf, text, *, size=14, bold=False, color=_BODY, bullet=False, space_before=2):
    p = tf.add_paragraph() if (tf.paragraphs and tf.paragraphs[0].runs) else tf.paragraphs[0]
    p.space_before = Pt(space_before)
    p.space_after = Pt(2)

    def _run(t, b):
        r = p.add_run()
        r.text = t
        r.font.size = Pt(size)
        r.font.bold = bold or b
        r.font.color.rgb = color

    if bullet:
        _run("•  ", False)
    # split on **bold** so emphasis survives into the deck (one run per segment)
    pos, segs = 0, []
    for m in _BOLD_RE.finditer(text):
        if m.start() > pos:
            segs.append((text[pos:m.start()], False))
        segs.append((m.group(1), True))
        pos = m.end()
    if pos < len(text):
        segs.append((text[pos:], False))
    for seg_text, seg_bold in (segs or [(text, False)]):
        if seg_text:
            _run(seg_text, seg_bold)
    return p


class _Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = _W
        self.prs.slide_height = _H
        self.section = None  # (slide, text_frame)

    def title_slide(self, doc: ExportDoc):
        s = _blank(self.prs)
        box = s.shapes.add_textbox(_MARGIN, Inches(2.3), _W - _MARGIN * 2, Inches(2.6))
        tf = box.text_frame
        tf.word_wrap = True
        _para(tf, doc.title, size=34, bold=True, color=_INK)
        if doc.subtitle and doc.subtitle != doc.title:
            _para(tf, doc.subtitle, size=15, color=_MUTED, space_before=10)
        if doc.meta:
            _para(tf, "   ·   ".join(doc.meta), size=11, color=_MUTED, space_before=14)
        bar = s.shapes.add_shape(1, _MARGIN, Inches(2.1), Inches(1.6), Pt(4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = _INDIGO
        bar.line.fill.background()

    def _new_section(self, heading: str):
        s = _blank(self.prs)
        _title_bar(s, heading)
        box = s.shapes.add_textbox(_MARGIN, Inches(1.7), _W - _MARGIN * 2, _H - Inches(2.4))
        tf = box.text_frame
        tf.word_wrap = True
        self.section = (s, tf)

    def _body_tf(self):
        if self.section is None:
            self._new_section("Details")
        return self.section[1]

    def image_slide(self, png: bytes, caption: str):
        s = _blank(self.prs)
        if caption:
            _title_bar(s, caption[:90])
        from PIL import Image as _PILImage
        iw, ih = _PILImage.open(io.BytesIO(png)).size
        max_w, max_h = _W - _MARGIN * 2, _H - Inches(2.0)
        w = max_w
        h = int(w * ih / iw)
        if h > max_h:
            h = max_h
            w = int(h * iw / ih)
        left = int((_W - w) / 2)
        s.shapes.add_picture(io.BytesIO(png), left, Inches(1.7), width=w, height=h)
        self.section = None

    def table_slide(self, columns, rows, caption: str):
        s = _blank(self.prs)
        _title_bar(s, caption or "Data")
        cols = columns[:6]
        body = rows[:12]
        n_r, n_c = len(body) + 1, len(cols)
        gtable = s.shapes.add_table(n_r, n_c, _MARGIN, Inches(1.7),
                                    _W - _MARGIN * 2, Inches(0.4) * n_r).table
        for j, c in enumerate(cols):
            cell = gtable.cell(0, j)
            cell.text = str(c)[:30]
            cell.fill.solid()
            cell.fill.fore_color.rgb = _INDIGO
            para = cell.text_frame.paragraphs[0]
            para.runs[0].font.size = Pt(11)
            para.runs[0].font.bold = True
            para.runs[0].font.color.rgb = _WHITE
        for i, r in enumerate(body, 1):
            for j in range(n_c):
                cell = gtable.cell(i, j)
                cell.text = (str(r[j])[:32] if j < len(r) else "")
                run = cell.text_frame.paragraphs[0].runs
                if run:
                    run[0].font.size = Pt(10)
                    run[0].font.color.rgb = _BODY
        self.section = None

    def add(self, b: Block):
        if b.kind == "heading":
            self._new_section(b.text)
        elif b.kind == "prose":
            tf = self._body_tf()
            if b.tag:
                _para(tf, b.tag.upper(), size=9, bold=True, color=_INDIGO, space_before=6)
            _para(tf, b.text, size=14, color=_BODY)
        elif b.kind == "bullets":
            tf = self._body_tf()
            for it in b.items:
                _para(tf, it, size=13, color=_BODY, bullet=True)
        elif b.kind == "finding":
            tf = self._body_tf()
            if b.caption:
                _para(tf, b.caption, size=14, bold=True, color=_INK, space_before=8)
            if b.text:
                _para(tf, b.text, size=12.5, color=_BODY)
            if b.confidence is not None:
                pct = int(round(b.confidence * 100))
                _para(tf, f"confidence: {pct}%", size=10, color=_MUTED)
        elif b.kind == "recs":
            tf = self._body_tf()
            for i, r in enumerate(b.recs, 1):
                _para(tf, f"{i}.  {r.get('action', '')}", size=13, color=_BODY, bullet=False, space_before=4)
                extra = [x for x in (
                    f"impact: {r['expected_impact']}" if r.get("expected_impact") else "",
                    f"owner: {r['owner']}" if r.get("owner") else "",
                    f"by {r['timeline']}" if r.get("timeline") else "",
                ) if x]
                if extra:
                    _para(tf, "    " + " · ".join(extra), size=10, color=_MUTED)
        elif b.kind == "keynums" and b.keynums:
            tf = self._body_tf()
            parts = [f"{k.label}: {k.value}" + (f" ({k.delta})" if k.delta else "") for k in b.keynums]
            _para(tf, "    ".join(parts), size=12.5, bold=True, color=_INK, space_before=4)
        elif b.kind == "chart" and b.png:
            self.image_slide(b.png, b.caption)
        elif b.kind == "table" and b.columns:
            self.table_slide(b.columns, b.rows, b.caption)
        # code blocks are omitted from slides (they live in the PDF appendix)

    def render(self) -> bytes:
        buf = io.BytesIO()
        self.prs.save(buf)
        buf.seek(0)
        return buf.read()


def render_pptx(doc: ExportDoc) -> bytes:
    deck = _Deck()
    deck.title_slide(doc)
    for b in doc.blocks:
        deck.add(b)
    return deck.render()
